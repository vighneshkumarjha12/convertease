from django.shortcuts import render, HttpResponse, redirect
from django.http import HttpResponse
from django.contrib import messages
from django.core.mail import send_mail
from django.conf import settings
from .models import Contact
from .forms import ContactForm
from PIL import Image
import io
import os
from io import BytesIO
# Heavy libraries are imported lazily inside each view to reduce startup memory usage


# ─────────────────────────────────────────────
#  Helper — keeps original name + appends _ce
# ─────────────────────────────────────────────
def get_output_filename(original_name, new_ext):
    """
    Example:
        get_output_filename('my_photo.jpg', 'png')  →  'my_photo_ce.png'
        get_output_filename('report.pdf',   'docx') →  'report_ce.docx'
    """
    base = os.path.splitext(original_name)[0]
    return f"{base}_ce.{new_ext}"


# ─────────────────────────────────────────────
#  Add Background
# ─────────────────────────────────────────────
def add_background(request):
    if request.method == 'POST':
        image    = request.FILES.get('image')
        bg_type  = request.POST.get('background_type')
        bg_color = request.POST.get('bg_color', '#FFFFFF')
        bg_image = request.FILES.get('bg_image')
        position = request.POST.get('position', 'center')

        original   = Image.open(image).convert("RGBA")

        if bg_type == 'color':
            background = Image.new('RGBA', original.size, bg_color)
        elif bg_type == 'image' and bg_image:
            background = Image.open(bg_image).convert("RGBA").resize(original.size)
        else:
            return HttpResponse("Invalid input", status=400)

        x, y = 0, 0
        if position == 'center':
            x = (background.width  - original.width)  // 2
            y = (background.height - original.height) // 2
        elif position == 'bottom-right':
            x = background.width  - original.width
            y = background.height - original.height

        background.paste(original, (x, y), original)

        buffer = io.BytesIO()
        background.save(buffer, format='PNG')
        buffer.seek(0)

        filename = get_output_filename(image.name, 'png')
        response = HttpResponse(buffer, content_type='image/png')
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response

    return render(request, 'add_background.html')


# ─────────────────────────────────────────────
#  Home
# ─────────────────────────────────────────────
def home(request):
    return render(request, 'home.html')


# ─────────────────────────────────────────────
#  PDF → PowerPoint
# ─────────────────────────────────────────────
def pdf_to_ppt(request):
    if request.method == 'POST':
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            pdf_file = request.FILES.get('pdf_file')
            if not pdf_file:
                return HttpResponse("No file uploaded.", status=400)

            pdf_data     = pdf_file.read()
            pdf_document = fitz.open(stream=pdf_data, filetype="pdf")
            presentation = Presentation()

            for page_number in range(len(pdf_document)):
                page   = pdf_document[page_number]
                text   = page.get_text("text")
                images = page.get_images(full=True)

                slide = presentation.slides.add_slide(presentation.slide_layouts[5])

                if text.strip():
                    textbox      = slide.shapes.add_textbox(Pt(10), Pt(10), Pt(960), Pt(540))
                    text_frame   = textbox.text_frame
                    p            = text_frame.add_paragraph()
                    p.text       = text
                    p.font.size  = Pt(14)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment  = PP_ALIGN.LEFT

                for img_index, img in enumerate(images):
                    xref         = img[0]
                    base_image   = pdf_document.extract_image(xref)
                    image_bytes  = base_image["image"]
                    image_ext    = base_image["ext"]

                    image_filename = f"temp_image_{page_number}_{img_index}.{image_ext}"
                    with open(image_filename, "wb") as img_file:
                        img_file.write(image_bytes)

                    slide.shapes.add_picture(image_filename, Pt(50), Pt(150), Pt(400), Pt(300))

            filename = get_output_filename(pdf_file.name, 'pptx')
            response = HttpResponse(
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            presentation.save(response)
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}", status=500)

    return render(request, 'pdf_to_ppt.html')


# ─────────────────────────────────────────────
#  Remove Background
# ─────────────────────────────────────────────
def remove_background_view(request):
    if request.method == "POST" and request.FILES.get("img"):
        from rembg import remove

        uploaded_file = request.FILES["img"]
        input_image   = Image.open(uploaded_file)

        output = remove(input_image)
        buffer = io.BytesIO()
        output.save(buffer, format="PNG")
        buffer.seek(0)

        filename = get_output_filename(uploaded_file.name, 'png')
        response = HttpResponse(buffer, content_type="image/png")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, 'remove_background.html')


# ─────────────────────────────────────────────
#  PDF → Word
# ─────────────────────────────────────────────
def pdf_to_word(request):
    if request.method == 'POST':
        try:
            import pdfplumber
            from docx import Document

            pdf_file = request.FILES.get('pdf_file')
            if not pdf_file:
                return HttpResponse("No file uploaded.", status=400)

            pdf_text = []
            with pdfplumber.open(pdf_file) as pdf:
                for page in pdf.pages:
                    pdf_text.append(page.extract_text())

            doc = Document()
            doc.add_heading("PDF to Word With ConvertEase", level=1)
            for page_num, page_text in enumerate(pdf_text, start=1):
                doc.add_heading(f"Page {page_num}", level=2)
                doc.add_paragraph(page_text)

            filename = get_output_filename(pdf_file.name, 'docx')
            response = HttpResponse(
                content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            doc.save(response)
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}", status=500)

    return render(request, 'pdf_to_word.html')


# ─────────────────────────────────────────────
#  Static pages
# ─────────────────────────────────────────────
def faq(request):
    return render(request, 'faq.html')

def about(request):
    return render(request, 'about.html')

def career(request):
    return render(request, 'career.html')

def htu(request):
    return render(request, 'htu.html')

def services(request):
    return render(request, 'services.html')

def blogs(request):
    return render(request, 'blogs.html')


# ─────────────────────────────────────────────
#  Image → PDF
# ─────────────────────────────────────────────
def imgtopdf(request):
    if request.method == 'POST':
        try:
            import img2pdf

            files = request.FILES.getlist('img')
            if not files:
                return HttpResponse("No files selected for conversion.")

            pdf      = img2pdf.convert([file.read() for file in files])
            filename = get_output_filename(files[0].name, 'pdf')
            response = HttpResponse(pdf, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}")

    return render(request, 'imgtopdf.html')


# ─────────────────────────────────────────────
#  Any Image → PNG
# ─────────────────────────────────────────────
def jpgtopng(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('img')
            if not files:
                return HttpResponse("No files selected.")

            buffer = io.BytesIO()
            im     = Image.open(files[0])
            im.save(buffer, "PNG")
            buffer.seek(0)

            filename = get_output_filename(files[0].name, 'png')
            response = HttpResponse(buffer, content_type="image/png")
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}")

    return render(request, 'jpgtopng.html')


# ─────────────────────────────────────────────
#  Any Image → JPG
# ─────────────────────────────────────────────
def pngtojpg(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('img')
            if not files:
                return HttpResponse("No files selected.")

            im = Image.open(files[0])
            if im.mode in ("RGBA", "LA"):
                background = Image.new("RGB", im.size, (255, 255, 255))
                background.paste(im, mask=im.split()[3])
                im = background
            else:
                im = im.convert('RGB')

            buffer = io.BytesIO()
            im.save(buffer, "JPEG")
            buffer.seek(0)

            filename = get_output_filename(files[0].name, 'jpg')
            response = HttpResponse(buffer, content_type="image/jpeg")
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}")

    return render(request, 'pngtojpg.html')


# ─────────────────────────────────────────────
#  Any Image → PNG  (webp route)
# ─────────────────────────────────────────────
def webptopng(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('img')
            if not files:
                return HttpResponse("No files selected.")

            buffer = io.BytesIO()
            im     = Image.open(files[0])
            im.save(buffer, "PNG")
            buffer.seek(0)

            filename = get_output_filename(files[0].name, 'png')
            response = HttpResponse(buffer, content_type="image/png")
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}")

    return render(request, 'webptopng.html')


# ─────────────────────────────────────────────
#  Any Image → WebP
# ─────────────────────────────────────────────
def pngtowebp(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('img')
            if not files:
                return HttpResponse("No files selected.")

            buffer = io.BytesIO()
            im     = Image.open(files[0])
            im.save(buffer, "WEBP")
            buffer.seek(0)

            filename = get_output_filename(files[0].name, 'webp')
            response = HttpResponse(buffer, content_type="image/webp")
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}")

    return render(request, 'pngtowebp.html')


# ─────────────────────────────────────────────
#  Contact
# ─────────────────────────────────────────────
def contact_view(request):
    if request.method == 'POST':
        form = ContactForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "Your message has been sent successfully!")
            return redirect('contact')
        else:
            messages.error(request, "There was an error in your form submission.")
    else:
        form = ContactForm()
    return render(request, 'contact.html', {'form': form})


# ─────────────────────────────────────────────
#  PDF → Excel
# ─────────────────────────────────────────────
def pdf_to_excel(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        try:
            import pandas as pd
            from PyPDF2 import PdfReader

            pdf_file   = request.FILES['pdf_file']
            pdf_reader = PdfReader(pdf_file)

            data = []
            for page in pdf_reader.pages:
                data.append(page.extract_text())

            df     = pd.DataFrame({'Content': data})
            output = BytesIO()
            writer = pd.ExcelWriter(output, engine='xlsxwriter')
            df.to_excel(writer, index=False, sheet_name='PDF Content')
            writer.close()

            filename = get_output_filename(pdf_file.name, 'xlsx')
            response = HttpResponse(
                output.getvalue(),
                content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            )
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except Exception as e:
            return HttpResponse(f"Error: {e}")

    return render(request, 'pdf_to_excel.html')